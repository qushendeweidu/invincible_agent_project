#!/usr/bin/env python3
"""PPTX 演示文稿生成工具。

根据结构化 JSON 输入生成 .pptx 文件，支持标题页、内容页、总结页。
被 presentation_agent 在收到演示文稿需求后调用。

依赖: pip install python-pptx

用法:
    python -m scripts.generate_pptx --input slides.json --output presentation.pptx
    echo '{"title": "...", "slides": [...]}' | python -m scripts.generate_pptx --stdin --output out.pptx
"""

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ====== 配色 ======
THEME = {
    "primary": RGBColor(0x63, 0x66, 0xF1) if HAS_PPTX else None,
    "dark": RGBColor(0x1E, 0x29, 0x3B) if HAS_PPTX else None,
    "light": RGBColor(0xF8, 0xFA, 0xFC) if HAS_PPTX else None,
    "gray": RGBColor(0x64, 0x74, 0x8B) if HAS_PPTX else None,
    "white": RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None,
}


def add_title_slide(prs: "Presentation", title: str, subtitle: str = ""):
    """添加标题页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["primary"]

    # 标题
    left, top, width, height = Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = THEME["white"]
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        left, top, width, height = Inches(1), Inches(4.2), Inches(8), Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = THEME["white"]
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs: "Presentation", title: str, bullets: list[str]):
    """添加内容页（标题 + 要点列表）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME["dark"]

    # 分隔线
    left, top = Inches(0.5), Inches(1.3)
    shape = slide.shapes.add_shape(1, left, top, Inches(9), Pt(2))  # 矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["primary"]
    shape.line.fill.background()

    # 要点
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = THEME["dark"]
        p.space_after = Pt(12)
        p.line_spacing = Pt(28)


def add_summary_slide(prs: "Presentation", title: str = "总结", points: list[str] = None):
    """添加总结页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["dark"]

    # 标题
    left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = THEME["white"]
    p.alignment = PP_ALIGN.CENTER

    if points:
        left, top, width, height = Inches(1), Inches(2.8), Inches(8), Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = THEME["white"]
            p.space_after = Pt(14)
            p.alignment = PP_ALIGN.CENTER


def generate(data: dict, output_path: str):
    """根据结构化数据生成 PPTX。"""
    if not HAS_PPTX:
        print("错误: 需要安装 python-pptx，执行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title = data.get("title", "演示文稿")
    subtitle = data.get("subtitle", "")
    slides = data.get("slides", [])
    summary = data.get("summary", [])

    # 标题页
    add_title_slide(prs, title, subtitle)

    # 内容页
    for slide_data in slides:
        slide_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        add_content_slide(prs, slide_title, bullets)

    # 总结页
    if summary:
        add_summary_slide(prs, "总结", summary)
    else:
        add_summary_slide(prs, "谢谢", ["感谢您的聆听"])

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"✅ 演示文稿已生成: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="PPTX 演示文稿生成器")
    parser.add_argument("--input", type=str, help="输入 JSON 文件路径")
    parser.add_argument("--stdin", action="store_true", help="从 stdin 读取 JSON")
    parser.add_argument("--output", type=str, default="output.pptx", help="输出 PPTX 路径")
    args = parser.parse_args()

    if args.stdin:
        data = json.load(sys.stdin)
    elif args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        # 演示用默认数据
        data = {
            "title": "项目演示",
            "subtitle": "由 AI 自动生成",
            "slides": [
                {"title": "项目背景", "bullets": ["行业现状分析", "市场需求", "核心痛点"]},
                {"title": "解决方案", "bullets": ["技术架构", "核心功能", "创新亮点"]},
                {"title": "实施计划", "bullets": ["第一阶段：需求分析", "第二阶段：开发实施", "第三阶段：测试上线"]},
            ],
            "summary": ["核心价值明确", "技术方案可行", "实施路径清晰"]
        }

    generate(data, args.output)


if __name__ == "__main__":
    main()
